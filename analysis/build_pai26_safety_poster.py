from pathlib import Path
from tempfile import NamedTemporaryFile
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path("/Users/abdulmohammad/Projects/Physics&Ling")
TEMPLATE = Path("/Users/abdulmohammad/Downloads/Copy of PAI26_Poster TEMPLATE.pptx")
OUT_DIR = BASE / "submission" / "PAI26_safety_variant" / "poster"
OUT = OUT_DIR / "PAI26_safety_poster.pptx"

GREEN = RGBColor(21, 95, 80)
TEAL = RGBColor(19, 106, 93)
PALE = RGBColor(225, 242, 240)
BLUE = RGBColor(42, 120, 189)
ORANGE = RGBColor(230, 126, 34)
MAGENTA = RGBColor(142, 35, 91)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(90, 90, 90)
LIGHT_GRAY = RGBColor(235, 239, 242)


def delete_shape(shape):
    shape._element.getparent().remove(shape._element)


def clear_text(shape):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    return tf


def set_textbox(shape, lines, font_size=22, color=BLACK, bold_first=False, italic_keywords=False):
    tf = clear_text(shape)
    for idx, item in enumerate(lines):
        text, size, bold, italic = item if isinstance(item, tuple) else (item, font_size, False, False)
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Georgia"
        p.font.bold = bold or (bold_first and idx == 0)
        p.font.italic = italic or (italic_keywords and idx == 1)
        p.space_after = Pt(3)


def add_box(slide, x, y, w, h, text, fill, line=WHITE, font_size=16, font_color=WHITE, bold=True):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(1.0)
    tf = clear_text(box)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Georgia"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = font_color
    return box


def add_label(slide, x, y, w, h, text, size=18, color=GREEN, bold=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = clear_text(tb)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Georgia"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return tb


def add_arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = GREEN
    line.line.width = Pt(2.2)
    line.line.end_arrowhead = True
    return line


def add_bar_chart(slide, x, y, w, h, title, categories, human_vals, ai_vals, max_val, y_label):
    add_label(slide, x, y, w, 0.35, title, size=17, color=GREEN, bold=True)
    plot_x = x + 0.25
    plot_y = y + 0.55
    plot_w = w - 0.45
    plot_h = h - 0.9

    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(plot_x), Inches(plot_y), Inches(plot_w), Inches(plot_h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(250, 252, 252)
    bg.line.color.rgb = LIGHT_GRAY
    bg.line.width = Pt(1)

    for frac in [0.25, 0.5, 0.75, 1.0]:
        yy = plot_y + plot_h - plot_h * frac
        grid = slide.shapes.add_connector(1, Inches(plot_x), Inches(yy), Inches(plot_x + plot_w), Inches(yy))
        grid.line.color.rgb = LIGHT_GRAY
        grid.line.width = Pt(0.7)

    n = len(categories)
    group_w = plot_w / n
    bar_w = min(0.22, group_w * 0.25)
    for i, cat in enumerate(categories):
        cx = plot_x + group_w * i + group_w / 2
        hv = human_vals[i]
        av = ai_vals[i]
        for offset, val, color in [(-bar_w * 0.7, hv, BLUE), (bar_w * 0.7, av, ORANGE)]:
            bh = plot_h * (val / max_val)
            bar = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(cx + offset - bar_w / 2),
                Inches(plot_y + plot_h - bh),
                Inches(bar_w),
                Inches(bh),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.color.rgb = color

        tb = slide.shapes.add_textbox(Inches(cx - group_w / 2 + 0.03), Inches(plot_y + plot_h + 0.08), Inches(group_w - 0.06), Inches(0.28))
        tf = clear_text(tb)
        p = tf.paragraphs[0]
        p.text = cat
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Georgia"
        p.font.size = Pt(9)
        p.font.color.rgb = BLACK

    add_label(slide, plot_x - 0.05, plot_y - 0.1, 1.2, 0.25, y_label, size=9, color=GRAY, bold=False)
    add_box(slide, x + w - 2.0, y + 0.02, 0.45, 0.22, "Human", BLUE, line=BLUE, font_size=7)
    add_box(slide, x + w - 1.05, y + 0.02, 0.45, 0.22, "AI", ORANGE, line=ORANGE, font_size=7)


def scrub_template_text(path: Path):
    replacements = {
        "TO DELETE!": "",
        "<Title for graphic / illustration / panel>": "",
        "<NAME>": "",
        "<add 5-10 keywords>": "",
        "<Scholar/Fellow/Student/Postdoc Name>": "",
        "<Dept Name>": "",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp_path = Path(tmp.name)
    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", compression=ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml"):
                text = data.decode("utf-8", errors="ignore")
                for old, new in replacements.items():
                    text = text.replace(old, new)
                data = text.encode("utf-8")
            zout.writestr(item, data)
    tmp_path.replace(path)


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(TEMPLATE))
    slide = prs.slides[0]
    shapes = list(slide.shapes)

    # Remove template-only graphic placeholders and orange instruction boxes.
    for idx in [13, 12, 2, 1]:
        delete_shape(shapes[idx])

    # Header and headline sections.
    set_textbox(
        shapes[7],
        [
            ("Abdul Mohammad", 31, True, False),
            ("English / Linguistics | AI safety, education evals, post-training", 13, False, False),
            ("Project: epistemic safety in AI physics tutors", 13, False, False),
            ("Open to research roles and internships in evals, alignment, or AI education", 12, False, False),
        ],
        color=WHITE,
    )
    set_textbox(
        shapes[6],
        [
            ("Can AI physics tutors be correct but incomplete?", 28, True, False),
            ("Keywords: epistemic safety; AI tutoring; physics education; LLM evaluation; discourse analysis", 15, False, True),
        ],
        color=WHITE,
    )
    set_textbox(
        shapes[8],
        [("contact", 17, True, False), ("email / LinkedIn", 11, False, False)],
        color=WHITE,
    )
    set_textbox(
        shapes[14],
        [("QR", 24, True, False), ("paper / repo", 13, False, False)],
        color=BLACK,
    )

    set_textbox(
        shapes[11],
        [
            ("Birds-eye view", 24, True, False),
            ("Correct answers can still be incomplete explanations for beginners.", 18, False, False),
            ("Risk: fluent responses can invite overtrust when reasoning steps are hidden.", 18, False, False),
            ("Dataset: 20 prompts, 60 explanations: 20 human and 40 AI.", 18, False, False),
            ("Method: sentence-level discourse labels plus 1-5 quality ratings.", 18, False, False),
            ("Reliability: kappa = 0.6706 on 72 overlap pairs.", 18, False, False),
            ("Baseline: keyword classifier reaches 70.6% accuracy and 0.575 macro-F1.", 18, False, False),
        ],
        color=BLACK,
    )
    add_box(slide, 0.65, 10.1, 4.65, 1.75, "Question\nCan a correct answer still be unsafe to learn from?", TEAL, font_size=13)
    add_box(slide, 5.65, 10.1, 4.65, 1.75, "Metric\nDoes the explanation expose reasoning structure?", TEAL, font_size=13)
    add_box(slide, 10.65, 10.1, 4.65, 1.75, "Finding\nAI uses fewer intuition-building bridges", MAGENTA, font_size=13)
    safety_card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(12.35), Inches(14.65), Inches(1.65))
    safety_card.fill.solid()
    safety_card.fill.fore_color.rgb = RGBColor(245, 250, 249)
    safety_card.line.color.rgb = GREEN
    safety_card.line.width = Pt(1.4)
    tf = clear_text(safety_card)
    p = tf.paragraphs[0]
    p.text = "Epistemic safety lens: fluent explanations can create confidence without giving students enough structure to detect gaps in understanding."
    p.font.name = "Georgia"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLACK

    set_textbox(
        shapes[9],
        [
            ("Implications for AI tutors", 23, True, False),
            ("Use discourse labels as deployment checks: missing framing, examples, caveats, or misconception repair.", 17, False, False),
            ("Use the same labels as post-training signals for explanations that teach reasoning, not just answers.", 17, False, False),
            ("Next: larger multi-rater corpus, unrevised human sources, and learner-outcome validation.", 17, False, False),
            ("Talk to me about evals, epistemic safety, education models, or post-training.", 17, False, False),
        ],
        color=BLACK,
    )

    # Right-side method panel.
    set_textbox(shapes[10], [("Study design: discourse labels as an evaluation layer", 23, True, False)], color=GREEN)
    add_box(slide, 18.0, 2.25, 4.0, 1.05, "20 intro physics prompts", TEAL)
    add_arrow(slide, 22.15, 2.78, 23.1, 2.78)
    add_box(slide, 23.25, 2.25, 4.2, 1.05, "Human + AI explanations", TEAL)
    add_arrow(slide, 27.6, 2.78, 28.55, 2.78)
    add_box(slide, 28.7, 2.25, 4.2, 1.05, "Sentence labels + ratings", TEAL)

    add_label(slide, 18.0, 3.85, 14.4, 0.45, "Five labels used to audit explanation structure", 18)
    labels = [
        ("FRAME", "sets up the physical situation"),
        ("PRINCIPLE", "states or applies a physics rule"),
        ("VERIFY", "checks a conclusion"),
        ("INTUITION", "gives an analogy or example"),
        ("CAVEAT", "marks assumptions or limits"),
    ]
    x0 = 18.0
    for i, (name, desc) in enumerate(labels):
        add_box(slide, x0 + i * 2.95, 4.45, 2.65, 0.85, f"{name}\n{desc}", MAGENTA if i in [3, 4] else TEAL, font_size=10)

    add_box(slide, 18.0, 6.0, 4.7, 1.1, "Reliability\nkappa = 0.6706", BLUE, font_size=15)
    add_box(slide, 23.3, 6.0, 4.7, 1.1, "Baseline\n70.6% accuracy", BLUE, font_size=15)
    add_box(slide, 28.6, 6.0, 4.2, 1.1, "Pilot scope\n60 explanations", BLUE, font_size=15)

    callout = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(18.0), Inches(8.0), Inches(14.8), Inches(2.15))
    callout.fill.solid()
    callout.fill.fore_color.rgb = PALE
    callout.line.color.rgb = RGBColor(190, 225, 220)
    tf = clear_text(callout)
    p = tf.paragraphs[0]
    p.text = "Safety framing: the target risk is misplaced learner confidence, not physical harm or hallucination. A fluent answer can be correct while omitting the conceptual bridge a beginner needs."
    p.font.name = "Georgia"
    p.font.size = Pt(17)
    p.font.color.rgb = BLACK

    # Right-side results panel.
    shapes[15].left = Inches(17.75)
    shapes[15].top = Inches(12.25)
    shapes[15].width = Inches(15.0)
    shapes[15].height = Inches(0.65)
    set_textbox(shapes[15], [("Key findings: correctness saturates; pedagogy separates", 23, True, False)], color=GREEN)
    add_bar_chart(
        slide,
        x=17.75,
        y=13.6,
        w=7.35,
        h=4.05,
        title="Move presence",
        categories=["FRAME", "PRINC.", "VERIFY", "INTUIT.", "CAVEAT"],
        human_vals=[100, 100, 65, 35, 40],
        ai_vals=[100, 100, 75, 20, 40],
        max_val=100,
        y_label="%",
    )
    add_bar_chart(
        slide,
        x=25.4,
        y=13.6,
        w=7.35,
        h=4.05,
        title="Quality ratings",
        categories=["clarity", "correct", "complete"],
        human_vals=[5.0, 5.0, 5.0],
        ai_vals=[4.875, 5.0, 4.0],
        max_val=5,
        y_label="1-5",
    )

    takeaway = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(17.75), Inches(18.35), Inches(15.0), Inches(2.35))
    takeaway.fill.solid()
    takeaway.fill.fore_color.rgb = RGBColor(245, 250, 249)
    takeaway.line.color.rgb = GREEN
    takeaway.line.width = Pt(1.5)
    tf = clear_text(takeaway)
    p = tf.paragraphs[0]
    p.text = "Takeaway: AI explanations matched human explanations on correctness, but had lower intuition-bridging coverage (20% vs 35%) and lower completeness ratings (4.0 vs 5.0)."
    p.font.name = "Georgia"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLACK

    prs.save(str(OUT))
    scrub_template_text(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
